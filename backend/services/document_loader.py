"""
StudyMate AI - Document Loader
Extracts raw text from PDFs, PPTX, DOCX, and TXT study materials.
"""

import os


def load_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    text = []
    for page in reader.pages:
        text.append(page.extract_text() or "")
    return "\n".join(text)


def load_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def load_docx(file_path: str) -> str:
    import docx
    doc = docx.Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def load_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def load_document(file_path: str) -> str:
    """Dispatch to the right loader based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext == ".pptx":
        return load_pptx(file_path)
    elif ext == ".docx":
        return load_docx(file_path)
    elif ext == ".txt":
        return load_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
